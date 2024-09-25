import os
from pptx import Presentation
from pptx.util import Inches,Pt
import FindMostSubDirs
import re

# Function to add 4 images to a slide
def add_images_to_slide(prs, img_paths):
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Define positions and sizes for 4 images with margins
    margin_top = Inches(0.5)
    margin_left = Inches(0.5)
    lr_margin_between = Inches(0.3)
    img_width = Inches(4)  # Image width
    img_height = Inches(3)  # Image height
    text_height = Inches(0.5)  # Space for text below each image

    # Define positions for 4 images
    # (left, top)
    positions = [(margin_left,margin_top), #top left
                 (lr_margin_between + img_width + margin_left, margin_top), #top right
                 (margin_left, img_height + margin_top+text_height), #bottom left
                 (lr_margin_between + img_width + margin_left, img_height + margin_top+text_height)] #bottom right

    # Add each image in the slide at the corresponding position
    #img_paths consists of 4 images
    #positions consists of 4 positions
    for img_path, (left, top) in zip(img_paths, positions):
        #add image
        slide.shapes.add_picture(img_path, left, top, img_width, img_height)

        #add image name
        # Get image name (file name without extension)
        img_name = os.path.basename(img_path).rsplit('.', 1)[0]

        # Add text box below the image with the image name
        text_box = slide.shapes.add_textbox(left, top + img_height, img_width, text_height)
        text_frame = text_box.text_frame
        text_frame.text = img_name

        # Optionally, customize text properties (font size, alignment)
        for paragraph in text_frame.paragraphs:
            paragraph.font.size = Pt(12)  # Font size for the image name

# Function to extract the numerical suffix from the image filename
def extract_suffix(filename):
    # Use regex to find the last number before the file extension (e.g., "01" in "3955_49_01.tif")
    match = re.search(r'_(\d+)\.tif$', filename)
    if match:
        return int(match.group(1))  # Convert the suffix to an integer for sorting
    return 0  # Default value if no match is found

#Program Start
root_directory = "/Users/adrianhwang/Documents/Xiaohui/New"
img_folders = FindMostSubDirs.get_most_subdirectories(root_directory) #list of directories with images

for img_folder in img_folders:

    # Create a PowerPoint presentation object
    prs = Presentation()

    # Get list of all images in the folder
    img_files = [f for f in os.listdir(img_folder) if f.endswith(('.tif'))]

    #sort the images in ascending order
    img_files.sort(key=extract_suffix)

    # Add images to slides, 4 per slide
    for i in range(0, len(img_files), 4):
        img_paths = [os.path.join(img_folder, img_files[j]) for j in range(i, min(i + 4, len(img_files)), 1)]
        add_images_to_slide(prs, img_paths)

    # Save the PowerPoint presentation
    basename = os.path.basename(img_folder).split(sep='-')[0].rstrip()
    ppt_name = basename + '.pptx'
    prs.save(os.path.join(img_folder, ppt_name))




