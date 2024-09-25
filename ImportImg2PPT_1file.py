import os
from pptx import Presentation
from pptx.util import Inches

# Set the directory containing images
img_folder = '/Users/adrianhwang/Documents/Xiaohui/Duvernay/Guidex Gvillee/2556.52-4.118_ TOC'

# Create a PowerPoint presentation object
prs = Presentation()

# Get list of all images in the folder
img_files = [f for f in os.listdir(img_folder) if f.endswith(('.tif'))]

# Function to add 4 images to a slide
def add_images_to_slide(prs, img_paths):
    slide_layout = prs.slide_layouts[5]  # Blank slide layout
    slide = prs.slides.add_slide(slide_layout)

    # Define positions and sizes for 4 images with margins
    margin_top = Inches(0.5)
    margin_left = Inches(0.5)
    margin_between = Inches(0.3)
    img_width = Inches(4)  # Image width
    img_height = Inches(3)  # Image height

    # Define positions for 4 images
    # (left, top)
    positions = [(margin_left,margin_top), #top left
                 (margin_between + img_width + margin_left, margin_top), #top right
                 (margin_left, margin_between + img_height + margin_top), #bottom left
                 (margin_between + img_width + margin_left, margin_between + img_height + margin_top)] #bottom right

    # Add each image in the slide at the corresponding position
    #img_paths consists of 4 images
    #positions consists of 4 positions
    for img_path, (left, top) in zip(img_paths, positions):
        slide.shapes.add_picture(img_path, left, top, img_width, img_height)

# Add images to slides, 4 per slide
for i in range(0, len(img_files), 4):
    img_paths = [os.path.join(img_folder, img_files[j]) for j in range(i, min(i + 4, len(img_files)),1)]
    add_images_to_slide(prs, img_paths)

# Save the PowerPoint presentation
basename=os.path.basename(img_folder).split(sep='-')[0].rstrip()
ppt_name=basename+'.pptx'
prs.save(os.path.join(img_folder, ppt_name))

